"""
Document parser — extracts raw text from PDF, DOCX, TXT, PPTX, and code files.
Always returns text for chunking (no raw-bytes path).
"""

import logging
import os
from pathlib import Path

logger = logging.getLogger(__name__)

try:
    from PyPDF2 import PdfReader
    PYPDF_AVAILABLE = True
except ImportError:
    PYPDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from tika import parser as tika_parser
    TIKA_AVAILABLE = True
except ImportError:
    TIKA_AVAILABLE = False


def extract_pdf(file_path: str) -> str:
    if not PYPDF_AVAILABLE:
        return ""
    try:
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        return ""


def extract_docx(file_path: str) -> str:
    if not DOCX_AVAILABLE:
        return ""
    try:
        doc = Document(file_path)
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text for c in row.cells))
        return "\n".join(parts)
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        return ""


def extract_pptx(file_path: str) -> str:
    if not PPTX_AVAILABLE:
        return ""
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        return ""


def extract_text_file(file_path: str) -> str:
    """Reads any text-based file (code, txt, md, etc.)"""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        logger.error(f"Text extraction failed: {e}")
        return ""


def extract_with_tika(file_path: str) -> str:
    if not TIKA_AVAILABLE:
        return ""
    try:
        raw = tika_parser.from_file(file_path)
        return raw.get("content", "") if raw else ""
    except Exception as e:
        logger.error(f"Tika failed: {e}")
        return ""


def parse_document(file_path: str) -> dict:
    """
    Extract text from any supported document or code file.
    Returns {"success": bool, "text": str, "error": str}
    """
    path = Path(file_path)

    if not path.exists():
        return {"success": False, "text": "", "error": "file not found"}

    ext = path.suffix.lower()
    text = ""

    # Specific parsers
    if ext == ".pdf":
        text = extract_pdf(file_path)
    elif ext in (".docx", ".doc"):
        text = extract_docx(file_path)
    elif ext in (".pptx", ".ppt"):
        text = extract_pptx(file_path)
    elif ext in (".txt", ".md", ".py", ".js", ".ts", ".c", ".cpp", ".h", ".html", ".css", ".json", ".yaml", ".yml", ".sh", ".sql"):
        text = extract_text_file(file_path)

    # Fallback to Tika if no text yet
    if not text or not text.strip():
        text = extract_with_tika(file_path)

    if not text or not text.strip():
        return {"success": False, "text": "", "error": "extraction failed"}

    return {"success": True, "text": text}
